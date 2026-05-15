from __future__ import annotations

from io import BytesIO
from pathlib import Path

import pypdf
import pytest
from docx import Document
from ebooklib import epub
from PIL import Image
from pptx import Presentation
from pypdf.generic import (
    ArrayObject,
    DecodedStreamObject,
    DictionaryObject,
    NameObject,
    NumberObject,
)

from services.api.app.operator_uploads import extractors


def _make_pdf(text: str) -> bytes:
    """Build a minimal one-page PDF containing the given Russian text."""
    writer = pypdf.PdfWriter()
    writer.add_blank_page(width=595, height=842)
    page = writer.pages[0]
    stream = DecodedStreamObject()
    content = (
        "BT /F1 12 Tf 72 770 Td (" + text.replace("(", "\\(").replace(")", "\\)") + ") Tj ET"
    )
    stream.set_data(content.encode("latin-1", errors="replace"))
    page[NameObject("/Contents")] = stream
    font = DictionaryObject(
        {
            NameObject("/Type"): NameObject("/Font"),
            NameObject("/Subtype"): NameObject("/Type1"),
            NameObject("/BaseFont"): NameObject("/Helvetica"),
        }
    )
    resources = DictionaryObject(
        {
            NameObject("/Font"): DictionaryObject({NameObject("/F1"): font}),
        }
    )
    page[NameObject("/Resources")] = resources
    page[NameObject("/MediaBox")] = ArrayObject(
        [NumberObject(0), NumberObject(0), NumberObject(595), NumberObject(842)]
    )
    buffer = BytesIO()
    writer.write(buffer)
    return buffer.getvalue()


def _write_docx(path: Path) -> None:
    document = Document()
    document.add_paragraph("Расписание работы офиса")
    document.add_paragraph("Понедельник — пятница: с 9:00 до 18:00")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Город"
    table.cell(0, 1).text = "Москва"
    table.cell(1, 0).text = "Телефон"
    table.cell(1, 1).text = "+7-495-000-00-00"
    document.save(str(path))


def _write_pptx(path: Path) -> None:
    presentation = Presentation()
    slide_layout = presentation.slide_layouts[5]
    slide = presentation.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Презентация для клиентов"
    textbox = slide.shapes.add_textbox(left=914400, top=914400, width=4572000, height=1828800)
    textbox.text_frame.text = "Услуги доступны 24/7"
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.text = "Не забыть упомянуть скидку"
    presentation.save(str(path))


def _write_png(path: Path) -> None:
    image = Image.new("RGB", (32, 32), color=(255, 255, 255))
    image.save(str(path), format="PNG")


def test_extract_pdf_returns_text(tmp_path: Path):
    pdf_path = tmp_path / "sample.pdf"
    pdf_path.write_bytes(_make_pdf("Hello PDF"))
    result = extractors.extract_pdf(pdf_path)
    assert "Hello PDF" in result


def test_extract_docx_returns_paragraphs_and_table_cells(tmp_path: Path):
    docx_path = tmp_path / "sample.docx"
    _write_docx(docx_path)
    result = extractors.extract_docx(docx_path)
    assert "Расписание работы офиса" in result
    assert "Понедельник" in result
    assert "Москва" in result
    assert "+7-495-000-00-00" in result


def test_extract_pptx_returns_slide_text_and_notes(tmp_path: Path):
    pptx_path = tmp_path / "sample.pptx"
    _write_pptx(pptx_path)
    result = extractors.extract_pptx(pptx_path)
    assert "Презентация для клиентов" in result
    assert "Услуги доступны 24/7" in result
    assert "Не забыть упомянуть скидку" in result


def test_extract_txt_decodes_utf8(tmp_path: Path):
    txt_path = tmp_path / "sample.txt"
    txt_path.write_text("Это тестовый файл", encoding="utf-8")
    result = extractors.extract_txt(txt_path)
    assert result == "Это тестовый файл"


def test_extract_txt_replaces_invalid_bytes(tmp_path: Path):
    txt_path = tmp_path / "broken.txt"
    txt_path.write_bytes(b"\xff\xfe" + "абв".encode("utf-8"))
    result = extractors.extract_txt(txt_path)
    assert "абв" in result


def test_extract_image_uses_pytesseract(monkeypatch, tmp_path: Path):
    png_path = tmp_path / "sample.png"
    _write_png(png_path)
    called_with: dict = {}

    def fake_ocr(image, lang):
        called_with["lang"] = lang
        return "Распознанный текст"

    monkeypatch.setattr(extractors.pytesseract, "image_to_string", fake_ocr)
    result = extractors.extract_image(png_path)
    assert "Распознанный текст" in result
    assert called_with["lang"] == "rus+eng"


def test_soft_wrap_preserves_short_sentences():
    text = "Короткое предложение. Второе короткое."
    result = extractors.soft_wrap(text, max_chars=100)
    assert result == "Короткое предложение.\nВторое короткое."


def test_soft_wrap_wraps_long_sentence_at_whitespace():
    sentence = "слово " * 60
    result = extractors.soft_wrap(sentence.strip(), max_chars=50)
    assert all(len(line) <= 50 for line in result.splitlines())
    assert "слово" in result


def test_soft_wrap_hard_cuts_when_no_whitespace():
    long_token = "x" * 250
    result = extractors.soft_wrap(long_token, max_chars=100)
    lines = result.splitlines()
    assert all(len(line) <= 100 for line in lines)
    assert "".join(lines) == long_token


def test_extract_dispatch_routes_to_extractor(tmp_path: Path):
    txt_path = tmp_path / "x.txt"
    txt_path.write_text("hello", encoding="utf-8")
    assert extractors.extract("txt", txt_path).strip() == "hello"


def test_extract_dispatch_raises_on_unknown_type(tmp_path: Path):
    with pytest.raises(extractors.ExtractionError) as exc:
        extractors.extract("bogus", tmp_path / "anything")
    assert "unsupported_file_type" in exc.value.reason


def test_extract_dispatch_raises_on_empty_text(tmp_path: Path):
    txt_path = tmp_path / "empty.txt"
    txt_path.write_text("   \n   \n", encoding="utf-8")
    with pytest.raises(extractors.ExtractionError) as exc:
        extractors.extract("txt", txt_path)
    assert exc.value.reason == "empty_text"


class FakeTranscriber:
    def __init__(self, return_text: str = "Расшифровка"):
        self.calls: list[tuple[Path, str]] = []
        self._return_text = return_text

    def transcribe(self, audio_path: Path, *, language: str) -> str:
        self.calls.append((audio_path, language))
        return self._return_text


@pytest.mark.asyncio
async def test_extract_audio_happy_path(monkeypatch, tmp_path: Path):
    audio = tmp_path / "voice.ogg"
    audio.write_bytes(b"opus-bytes")
    monkeypatch.setattr(extractors, "_probe_duration", lambda path: 12.5)
    transcriber = FakeTranscriber("Привет это тестовое сообщение")
    result = await extractors.extract_audio(audio, transcriber=transcriber, max_seconds=60)
    assert "Привет" in result
    assert transcriber.calls[0][1] == "ru"


@pytest.mark.asyncio
async def test_extract_audio_too_long_skips_transcription(monkeypatch, tmp_path: Path):
    audio = tmp_path / "long.mp3"
    audio.write_bytes(b"x")
    monkeypatch.setattr(extractors, "_probe_duration", lambda path: 9999.0)
    transcriber = FakeTranscriber()
    with pytest.raises(extractors.ExtractionError) as exc:
        await extractors.extract_audio(audio, transcriber=transcriber, max_seconds=60)
    assert exc.value.reason == "audio_too_long"
    assert transcriber.calls == []


@pytest.mark.asyncio
async def test_extract_audio_empty_transcript_raises(monkeypatch, tmp_path: Path):
    audio = tmp_path / "silence.mp3"
    audio.write_bytes(b"x")
    monkeypatch.setattr(extractors, "_probe_duration", lambda path: 1.0)
    transcriber = FakeTranscriber(return_text="   ")
    with pytest.raises(extractors.ExtractionError) as exc:
        await extractors.extract_audio(audio, transcriber=transcriber, max_seconds=60)
    assert exc.value.reason == "empty_text"


@pytest.mark.asyncio
async def test_extract_audio_uses_settings_default_cap(monkeypatch, tmp_path: Path):
    audio = tmp_path / "v.mp3"
    audio.write_bytes(b"x")
    monkeypatch.setattr(extractors, "_probe_duration", lambda path: 5.0)
    transcriber = FakeTranscriber("ok")
    result = await extractors.extract_audio(audio, transcriber=transcriber)
    assert result == "ok"


@pytest.mark.asyncio
async def test_extract_video_runs_ffmpeg_then_audio(monkeypatch, tmp_path: Path):
    video = tmp_path / "v.mp4"
    video.write_bytes(b"video-bytes")
    monkeypatch.setattr(extractors, "_probe_duration", lambda path: 30.0)

    captured_cmds: list[list[str]] = []

    def fake_run(cmd, **kwargs):
        captured_cmds.append(cmd)
        Path(cmd[-1]).write_bytes(b"audio-data")
        return subprocess_completed(cmd)

    monkeypatch.setattr(extractors.subprocess, "run", fake_run)
    transcriber = FakeTranscriber("video transcript")
    result = await extractors.extract_video(video, transcriber=transcriber, max_seconds=120)
    assert "video transcript" in result
    assert captured_cmds[0][0] == "ffmpeg"


@pytest.mark.asyncio
async def test_extract_video_too_long_skips_ffmpeg(monkeypatch, tmp_path: Path):
    video = tmp_path / "long.mp4"
    video.write_bytes(b"x")
    monkeypatch.setattr(extractors, "_probe_duration", lambda path: 99999.0)
    called = {"ran": False}

    def fake_run(*args, **kwargs):
        called["ran"] = True
        return subprocess_completed(args[0])

    monkeypatch.setattr(extractors.subprocess, "run", fake_run)
    with pytest.raises(extractors.ExtractionError) as exc:
        await extractors.extract_video(video, transcriber=FakeTranscriber(), max_seconds=60)
    assert exc.value.reason == "audio_too_long"
    assert called["ran"] is False


def subprocess_completed(cmd):
    import subprocess as _sp

    return _sp.CompletedProcess(args=cmd, returncode=0, stdout="", stderr="")


def test_probe_duration_parses_ffprobe_output(monkeypatch, tmp_path: Path):
    media = tmp_path / "m.mp3"
    media.write_bytes(b"x")

    def fake_run(cmd, **kwargs):
        return subprocess_completed_with_stdout(cmd, "42.5\n")

    monkeypatch.setattr(extractors.subprocess, "run", fake_run)
    assert extractors._probe_duration(media) == 42.5


def test_probe_duration_raises_on_empty_output(monkeypatch, tmp_path: Path):
    media = tmp_path / "m.mp3"
    media.write_bytes(b"x")
    monkeypatch.setattr(
        extractors.subprocess,
        "run",
        lambda cmd, **kwargs: subprocess_completed_with_stdout(cmd, "   \n"),
    )
    with pytest.raises(extractors.ExtractionError) as exc:
        extractors._probe_duration(media)
    assert exc.value.reason == "ffprobe_no_duration"


def subprocess_completed_with_stdout(cmd, stdout):
    import subprocess as _sp

    return _sp.CompletedProcess(args=cmd, returncode=0, stdout=stdout, stderr="")


@pytest.mark.asyncio
async def test_extract_media_dispatches_audio_and_video(monkeypatch, tmp_path: Path):
    monkeypatch.setattr(extractors, "_probe_duration", lambda path: 5.0)
    audio = tmp_path / "a.mp3"
    audio.write_bytes(b"x")
    video = tmp_path / "v.mp4"
    video.write_bytes(b"x")

    async def stub_extract_video(path, *, transcriber, max_seconds=None):
        return "video!"

    monkeypatch.setattr(extractors, "extract_video", stub_extract_video)
    result_audio = await extractors.extract_media(
        "audio", audio, transcriber=FakeTranscriber("audio!")
    )
    result_video = await extractors.extract_media(
        "video", video, transcriber=FakeTranscriber("ignored")
    )
    assert result_audio == "audio!"
    assert result_video == "video!"


@pytest.mark.asyncio
async def test_extract_media_rejects_unknown_type(tmp_path: Path):
    with pytest.raises(extractors.ExtractionError) as exc:
        await extractors.extract_media(
            "pdf", tmp_path / "x", transcriber=FakeTranscriber()
        )
    assert "unsupported_media_type" in exc.value.reason


def test_binary_sha256_streams(tmp_path: Path):
    path = tmp_path / "blob.bin"
    contents = b"A" * 200_000
    path.write_bytes(contents)
    import hashlib

    expected = hashlib.sha256(contents).hexdigest()
    assert extractors.binary_sha256(path) == expected


def test_whisper_transcriber_uses_lazy_loaded_model(monkeypatch):
    transcriber = extractors.WhisperTranscriber()

    class FakeSegment:
        def __init__(self, text: str) -> None:
            self.text = text

    class FakeModel:
        def __init__(self) -> None:
            self.calls: list[tuple[str, str]] = []

        def transcribe(self, audio_path: str, *, language: str):
            self.calls.append((audio_path, language))
            return ([FakeSegment("hello"), FakeSegment("world")], None)

    fake_model = FakeModel()
    monkeypatch.setattr(transcriber, "_load", lambda: fake_model)
    result = transcriber.transcribe(Path("/tmp/x.mp3"), language="ru")
    assert result == "hello world"
    assert fake_model.calls[0] == ("/tmp/x.mp3", "ru")
    # Second call reuses the same model
    transcriber.transcribe(Path("/tmp/y.mp3"), language="ru")
    assert len(fake_model.calls) == 2


def test_extractors_registry_includes_new_formats():
    for key in ("xlsx", "csv", "html", "md", "rtf", "epub", "zip"):
        assert key in extractors.EXTRACTORS


def _write_xlsx(path: Path) -> None:
    from openpyxl import Workbook

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Контакты"
    sheet.append(["Город", "Телефон"])
    sheet.append(["Москва", "+7-495-000-00-00"])
    sheet.append([None, None])  # blank row should be skipped
    second = workbook.create_sheet("Расписание")
    second.append(["Понедельник", "9:00-18:00"])
    workbook.save(str(path))


def test_extract_xlsx_emits_sheet_headers_and_rows(tmp_path: Path):
    xlsx_path = tmp_path / "sample.xlsx"
    _write_xlsx(xlsx_path)
    result = extractors.extract_xlsx(xlsx_path)
    assert "# Sheet: Контакты" in result
    assert "Москва\t+7-495-000-00-00" in result
    assert "# Sheet: Расписание" in result
    assert "Понедельник\t9:00-18:00" in result


def test_extract_csv_handles_comma_and_pipe_join(tmp_path: Path):
    csv_path = tmp_path / "table.csv"
    csv_path.write_text("город,телефон\nМосква,+7-495-000\n", encoding="utf-8")
    result = extractors.extract_csv(csv_path)
    assert "город | телефон" in result
    assert "Москва | +7-495-000" in result


def test_extract_csv_handles_semicolon_delimiter(tmp_path: Path):
    csv_path = tmp_path / "semi.csv"
    csv_path.write_text("a;b;c\n1;2;3\n", encoding="utf-8")
    result = extractors.extract_csv(csv_path)
    assert "a | b | c" in result
    assert "1 | 2 | 3" in result


def test_extract_csv_falls_back_to_excel_dialect_on_sniff_error(tmp_path: Path):
    csv_path = tmp_path / "tricky.csv"
    csv_path.write_text("solo\n", encoding="utf-8")
    result = extractors.extract_csv(csv_path)
    assert result == "solo"


def test_extract_csv_returns_empty_string_for_whitespace_only(tmp_path: Path):
    csv_path = tmp_path / "blank.csv"
    csv_path.write_text("   \n\n   ", encoding="utf-8")
    assert extractors.extract_csv(csv_path) == ""


def test_extract_html_strips_scripts_and_styles(tmp_path: Path):
    html_path = tmp_path / "page.html"
    html_path.write_text(
        "<html><head><style>body{}</style><script>alert('xss')</script>"
        "</head><body><p>Текст страницы</p><noscript>hidden</noscript></body></html>",
        encoding="utf-8",
    )
    result = extractors.extract_html(html_path)
    assert "Текст страницы" in result
    assert "alert" not in result
    assert "hidden" not in result


def test_extract_md_strips_headings_bullets_and_links(tmp_path: Path):
    md_path = tmp_path / "doc.md"
    md_path.write_text(
        "# Заголовок\n\n- пункт **один**\n- пункт _два_\n\n"
        "Ссылка [сюда](https://example.com) и `код`.\n\n"
        "```python\nprint('x')\n```\n",
        encoding="utf-8",
    )
    result = extractors.extract_md(md_path)
    assert "Заголовок" in result
    assert "пункт один" in result
    assert "пункт два" in result
    assert "Ссылка сюда и код." in result
    assert "https://example.com" not in result
    assert "```" not in result


def test_extract_rtf_returns_plain_text(tmp_path: Path):
    rtf_path = tmp_path / "note.rtf"
    rtf_path.write_text(
        r"{\rtf1\ansi\ansicpg1251 Привет мир из RTF.}", encoding="utf-8"
    )
    result = extractors.extract_rtf(rtf_path)
    assert "Привет мир из RTF" in result


def _write_epub(path: Path) -> None:
    book = epub.EpubBook()
    book.set_identifier("test-id")
    book.set_title("Тестовая книга")
    book.set_language("ru")
    chap1 = epub.EpubHtml(title="Глава 1", file_name="c1.xhtml", lang="ru")
    chap1.content = (
        "<html><body><h1>Глава 1</h1><p>Первая глава.</p></body></html>"
    )
    chap2 = epub.EpubHtml(title="Глава 2", file_name="c2.xhtml", lang="ru")
    chap2.content = (
        "<html><body><h1>Глава 2</h1><p>Вторая глава.</p></body></html>"
    )
    book.add_item(chap1)
    book.add_item(chap2)
    book.toc = (chap1, chap2)
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    book.spine = ["nav", chap1, chap2]
    epub.write_epub(str(path), book)


def test_extract_epub_concatenates_chapters(tmp_path: Path):
    epub_path = tmp_path / "book.epub"
    _write_epub(epub_path)
    result = extractors.extract_epub(epub_path)
    assert "Первая глава." in result
    assert "Вторая глава." in result


def test_extract_zip_dispatches_to_member_extractors(tmp_path: Path):
    import zipfile as zf

    zip_path = tmp_path / "bundle.zip"
    with zf.ZipFile(zip_path, "w") as archive:
        archive.writestr("note.txt", "Содержимое заметки".encode("utf-8"))
        archive.writestr("data.csv", "город,телефон\nМосква,+7-495\n".encode("utf-8"))
        archive.writestr("ignore.exe", b"binary")  # unsupported, must be skipped
    result = extractors.extract_zip(zip_path)
    assert "--- note.txt ---" in result
    assert "Содержимое заметки" in result
    assert "--- data.csv ---" in result
    assert "Москва | +7-495" in result
    assert "binary" not in result


def test_extract_zip_rejects_nested_archive(tmp_path: Path):
    import zipfile as zf

    zip_path = tmp_path / "outer.zip"
    with zf.ZipFile(zip_path, "w") as archive:
        archive.writestr("inner.zip", b"PK")
    with pytest.raises(extractors.ExtractionError) as exc:
        extractors.extract_zip(zip_path)
    assert exc.value.reason == "nested_zip_not_supported"


def test_extract_zip_rejects_too_many_members(monkeypatch, tmp_path: Path):
    import zipfile as zf

    zip_path = tmp_path / "many.zip"
    monkeypatch.setattr(extractors, "_ZIP_MAX_MEMBERS", 2)
    with zf.ZipFile(zip_path, "w") as archive:
        for i in range(3):
            archive.writestr(f"f{i}.txt", b"x")
    with pytest.raises(extractors.ExtractionError) as exc:
        extractors.extract_zip(zip_path)
    assert exc.value.reason == "zip_too_many_members"


def test_extract_zip_rejects_total_uncompressed_over_cap(monkeypatch, tmp_path: Path):
    import zipfile as zf

    zip_path = tmp_path / "big.zip"
    monkeypatch.setattr(extractors, "_ZIP_MAX_TOTAL_UNCOMPRESSED", 10)
    with zf.ZipFile(zip_path, "w") as archive:
        archive.writestr("a.txt", b"x" * 20)
    with pytest.raises(extractors.ExtractionError) as exc:
        extractors.extract_zip(zip_path)
    assert exc.value.reason == "zip_too_large"


def test_extract_zip_skips_member_over_individual_size_cap(monkeypatch, tmp_path: Path):
    import zipfile as zf

    from platform_common import settings as settings_module

    zip_path = tmp_path / "mixed.zip"

    class FakeSettings:
        operator_upload_max_bytes = 5

    monkeypatch.setattr(settings_module, "get_settings", lambda: FakeSettings())
    monkeypatch.setattr(extractors, "get_settings", lambda: FakeSettings())
    with zf.ZipFile(zip_path, "w") as archive:
        archive.writestr("small.txt", b"hi")
        archive.writestr("huge.txt", b"x" * 1000)
    result = extractors.extract_zip(zip_path)
    assert "--- small.txt ---" in result
    assert "--- huge.txt ---" not in result


def test_extract_zip_skips_member_extractor_failures(monkeypatch, tmp_path: Path):
    import zipfile as zf

    zip_path = tmp_path / "mixed.zip"
    with zf.ZipFile(zip_path, "w") as archive:
        archive.writestr("ok.txt", b"hello")
        archive.writestr("blank.txt", b"   ")
    result = extractors.extract_zip(zip_path)
    assert "--- ok.txt ---" in result
    assert "hello" in result
    assert "--- blank.txt ---" not in result


def test_extract_zip_skips_member_raising_extraction_error(monkeypatch, tmp_path: Path):
    import zipfile as zf

    zip_path = tmp_path / "mixed.zip"
    with zf.ZipFile(zip_path, "w") as archive:
        archive.writestr("ok.txt", b"good")
        archive.writestr("flaky.txt", b"bad")
    real_txt = extractors.EXTRACTORS["txt"]
    invocations: list[Path] = []

    def flaky(path: Path) -> str:
        invocations.append(path)
        if path.name == "member.txt" and len(invocations) == 2:
            raise extractors.ExtractionError("synthetic_failure")
        return real_txt(path)

    monkeypatch.setitem(extractors.EXTRACTORS, "txt", flaky)
    result = extractors.extract_zip(zip_path)
    assert "good" in result
    assert "bad" not in result


def test_extract_zip_bad_zipfile_raises(tmp_path: Path):
    zip_path = tmp_path / "garbage.zip"
    zip_path.write_bytes(b"not a zip at all")
    with pytest.raises(extractors.ExtractionError) as exc:
        extractors.extract_zip(zip_path)
    assert exc.value.reason == "zip_corrupt"


def test_extract_zip_testzip_reports_bad_member_raises(monkeypatch, tmp_path: Path):
    import zipfile as zf

    zip_path = tmp_path / "ok.zip"
    with zf.ZipFile(zip_path, "w") as archive:
        archive.writestr("a.txt", b"hello")
    monkeypatch.setattr(zf.ZipFile, "testzip", lambda self: "a.txt")
    with pytest.raises(extractors.ExtractionError) as exc:
        extractors.extract_zip(zip_path)
    assert exc.value.reason == "zip_corrupt"


@pytest.mark.parametrize(
    "source_type,filename,payload",
    [
        ("csv", "blank.csv", "   \n\n"),
        ("html", "blank.html", "<html><body></body></html>"),
        ("md", "blank.md", "\n\n"),
        ("rtf", "blank.rtf", r"{\rtf1 }"),
    ],
)
def test_extract_dispatch_raises_empty_text_for_new_text_formats(
    tmp_path: Path, source_type: str, filename: str, payload: str
):
    file_path = tmp_path / filename
    file_path.write_text(payload, encoding="utf-8")
    with pytest.raises(extractors.ExtractionError) as exc:
        extractors.extract(source_type, file_path)
    assert exc.value.reason == "empty_text"


def test_extract_dispatch_raises_empty_text_for_empty_zip(tmp_path: Path):
    import zipfile as zf

    zip_path = tmp_path / "empty.zip"
    with zf.ZipFile(zip_path, "w"):
        pass
    with pytest.raises(extractors.ExtractionError) as exc:
        extractors.extract("zip", zip_path)
    assert exc.value.reason == "empty_text"


def test_extract_epub_skips_chapters_with_empty_extracted_text(tmp_path: Path):
    epub_path = tmp_path / "mixed.epub"
    book = epub.EpubBook()
    book.set_identifier("id")
    book.set_title("t")
    book.set_language("ru")
    real = epub.EpubHtml(title="real", file_name="real.xhtml", lang="ru")
    real.content = "<html><body><p>Реальный текст.</p></body></html>"
    blank = epub.EpubHtml(title="blank", file_name="blank.xhtml", lang="ru")
    blank.content = (
        "<html><body><script>console.log('hidden')</script></body></html>"
    )
    book.add_item(real)
    book.add_item(blank)
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    book.spine = ["nav", real, blank]
    epub.write_epub(str(epub_path), book)
    result = extractors.extract_epub(epub_path)
    assert "Реальный текст" in result
    assert "console.log" not in result
    assert "hidden" not in result
