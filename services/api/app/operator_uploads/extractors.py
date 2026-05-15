"""Local file extractors for operator KB uploads.

Every extractor returns the raw text content of the source file. Empty
output (no usable text) raises `ExtractionError("empty_text")`. All
extractors run on CPU with zero external API calls — PDF, DOCX, PPTX,
XLSX, CSV, HTML, Markdown, RTF, EPUB and TXT use pure-Python libraries;
images use the local `tesseract` binary through `pytesseract` (Russian +
English language packs installed at container build time); audio/video
use `faster-whisper` locally with an upper duration cap. ZIP archives
are walked one level deep and each supported member is dispatched back
through this module.
"""

from __future__ import annotations

import asyncio
import csv as csv_module
import hashlib
import io
import re
import subprocess
import tempfile
import zipfile
from pathlib import Path
from typing import Protocol

import pypdf
import pytesseract
from bs4 import BeautifulSoup
from docx import Document
from ebooklib import ITEM_DOCUMENT, epub
from openpyxl import load_workbook
from PIL import Image
from pptx import Presentation
from striprtf.striprtf import rtf_to_text

from platform_common.settings import get_settings
from services.api.app.russian_text import get_russian_normalizer


class ExtractionError(Exception):
    def __init__(self, reason: str) -> None:
        super().__init__(reason)
        self.reason = reason


def extract_pdf(path: Path) -> str:
    reader = pypdf.PdfReader(str(path))
    pages: list[str] = []
    for page in reader.pages:
        text = page.extract_text() or ""
        text = text.strip()
        if text:
            pages.append(text)
    return "\n\n".join(pages)


def extract_docx(path: Path) -> str:
    document = Document(str(path))
    parts: list[str] = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)
    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                text = cell.text.strip()
                if text:
                    parts.append(text)
    return "\n".join(parts)


def extract_pptx(path: Path) -> str:
    presentation = Presentation(str(path))
    parts: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in paragraph.runs).strip()
                    if text:
                        parts.append(text)
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame is not None:
                notes_text = notes_frame.text.strip()
                if notes_text:
                    parts.append(notes_text)
    return "\n".join(parts)


def extract_txt(path: Path) -> str:
    return path.read_bytes().decode("utf-8", errors="replace")


def extract_image(path: Path) -> str:
    with Image.open(str(path)) as image:
        return pytesseract.image_to_string(image, lang="rus+eng")


def extract_xlsx(path: Path) -> str:
    workbook = load_workbook(filename=str(path), read_only=True, data_only=True)
    parts: list[str] = []
    try:
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            parts.append(f"# Sheet: {sheet_name}")
            for row in sheet.iter_rows(values_only=True):
                cells = [str(cell) for cell in row if cell is not None and str(cell).strip()]
                if cells:
                    parts.append("\t".join(cells))
    finally:
        workbook.close()
    return "\n".join(parts)


def extract_csv(path: Path) -> str:
    raw = path.read_bytes().decode("utf-8", errors="replace")
    if not raw.strip():
        return ""
    sample = raw[:4096]
    try:
        dialect = csv_module.Sniffer().sniff(sample, delimiters=",;\t|")
    except csv_module.Error:
        dialect = csv_module.excel
    reader = csv_module.reader(io.StringIO(raw), dialect)
    parts: list[str] = []
    for row in reader:
        cells = [cell.strip() for cell in row if cell and cell.strip()]
        if cells:
            parts.append(" | ".join(cells))
    return "\n".join(parts)


def _html_to_text(source: str | bytes) -> str:
    """Strip script/style/template/noscript subtrees and return visible text."""
    soup = BeautifulSoup(source, "html.parser")
    for tag in soup(["script", "style", "noscript", "template"]):
        tag.decompose()
    return soup.get_text(separator="\n", strip=True)


def extract_html(path: Path) -> str:
    return _html_to_text(path.read_bytes())


_MD_FENCE_RE = re.compile(r"^```.*$", re.MULTILINE)
_MD_PREFIX_RE = re.compile(r"^\s*(?:#{1,6}\s+|>+\s*|[-*+]\s+)", re.MULTILINE)
_MD_LINK_RE = re.compile(r"\[([^\]]+)\]\([^)]+\)")
_MD_EMPHASIS_RE = re.compile(r"(\*\*|__|\*|_)(.+?)\1")
_MD_INLINE_CODE_RE = re.compile(r"`([^`]+)`")


def extract_md(path: Path) -> str:
    raw = path.read_bytes().decode("utf-8", errors="replace")
    text = _MD_FENCE_RE.sub("", raw)
    text = _MD_LINK_RE.sub(r"\1", text)
    text = _MD_EMPHASIS_RE.sub(r"\2", text)
    text = _MD_INLINE_CODE_RE.sub(r"\1", text)
    text = _MD_PREFIX_RE.sub("", text)
    lines = [line.strip() for line in text.splitlines()]
    return "\n".join(line for line in lines if line)


def extract_rtf(path: Path) -> str:
    raw = path.read_bytes().decode("utf-8", errors="replace")
    return rtf_to_text(raw, errors="ignore")


def extract_epub(path: Path) -> str:
    book = epub.read_epub(str(path))
    parts: list[str] = []
    for item in book.get_items_of_type(ITEM_DOCUMENT):
        text = _html_to_text(item.get_content())
        if text.strip():
            parts.append(text)
    return "\n\n".join(parts)


_ZIP_MAX_MEMBERS = 100
_ZIP_MAX_TOTAL_UNCOMPRESSED = 50 * 1024 * 1024
_ZIP_MEMBER_EXTENSIONS: dict[str, str] = {
    ".pdf": "pdf",
    ".docx": "docx",
    ".pptx": "pptx",
    ".txt": "txt",
    ".png": "image",
    ".jpg": "image",
    ".jpeg": "image",
    ".gif": "image",
    ".bmp": "image",
    ".webp": "image",
    ".tiff": "image",
    ".xlsx": "xlsx",
    ".csv": "csv",
    ".html": "html",
    ".htm": "html",
    ".md": "md",
    ".markdown": "md",
    ".rtf": "rtf",
    ".epub": "epub",
}


def extract_zip(path: Path) -> str:
    try:
        archive = zipfile.ZipFile(str(path))
    except zipfile.BadZipFile as exc:
        raise ExtractionError("zip_corrupt") from exc
    with archive:
        if archive.testzip() is not None:
            raise ExtractionError("zip_corrupt")
        members = [info for info in archive.infolist() if not info.is_dir()]
        if len(members) > _ZIP_MAX_MEMBERS:
            raise ExtractionError("zip_too_many_members")
        total = sum(info.file_size for info in members)
        if total > _ZIP_MAX_TOTAL_UNCOMPRESSED:
            raise ExtractionError("zip_too_large")
        max_member_bytes = get_settings().operator_upload_max_bytes
        parts: list[str] = []
        for info in members:
            member_name = info.filename
            suffix = Path(member_name).suffix.lower()
            if suffix == ".zip":
                raise ExtractionError("nested_zip_not_supported")
            extractor_key = _ZIP_MEMBER_EXTENSIONS.get(suffix)
            if extractor_key is None:
                continue
            if info.file_size > max_member_bytes:
                continue
            extractor = EXTRACTORS[extractor_key]
            with tempfile.TemporaryDirectory() as tmp_dir:
                target = Path(tmp_dir) / f"member{suffix}"
                with archive.open(info) as src, target.open("wb") as dst:
                    dst.write(src.read())
                try:
                    text = extractor(target)
                except ExtractionError:
                    continue
            if text and text.strip():
                parts.append(f"--- {member_name} ---\n{text}")
    return "\n\n".join(parts)


def soft_wrap(text: str, *, max_chars: int = 200) -> str:
    """Sentence-segment text and wrap lines softly at `max_chars`.

    The RAG ingest pipeline uses `text.splitlines()` to produce chunks
    (one per non-empty line). Without segmentation, a long extracted
    paragraph would become a single multi-thousand-char chunk that
    retrieval scores poorly. We:

    1. Split into Russian sentences via razdel (`RussianNormalizer.sentenize`).
    2. For any sentence still longer than `max_chars`, break it at the
       last whitespace before `max_chars` so chunks stay retrievable.

    Lines emitted are non-empty; the result joins them with `\\n`.
    """
    normalizer = get_russian_normalizer()
    sentences = normalizer.sentenize(text)
    out: list[str] = []
    for sentence in sentences:
        remaining = sentence
        while len(remaining) > max_chars:
            cutoff = remaining.rfind(" ", 0, max_chars)
            if cutoff <= 0:
                cutoff = max_chars
            piece = remaining[:cutoff].strip()
            if piece:
                out.append(piece)
            remaining = remaining[cutoff:].strip()
        if remaining:
            out.append(remaining)
    return "\n".join(out)


EXTRACTORS: dict[str, callable] = {
    "pdf": extract_pdf,
    "docx": extract_docx,
    "pptx": extract_pptx,
    "txt": extract_txt,
    "image": extract_image,
    "xlsx": extract_xlsx,
    "csv": extract_csv,
    "html": extract_html,
    "md": extract_md,
    "rtf": extract_rtf,
    "epub": extract_epub,
    "zip": extract_zip,
}


def extract(source_file_type: str, path: Path) -> str:
    """Dispatch to the right extractor and raise on empty output."""
    extractor = EXTRACTORS.get(source_file_type)
    if extractor is None:
        raise ExtractionError(f"unsupported_file_type:{source_file_type}")
    raw = extractor(path)
    if not raw or not raw.strip():
        raise ExtractionError("empty_text")
    return raw


class Transcriber(Protocol):
    def transcribe(self, audio_path: Path, *, language: str) -> str: ...


def _probe_duration(path: Path) -> float:
    """Return media duration in seconds via ffprobe."""
    completed = subprocess.run(
        [
            "ffprobe",
            "-v",
            "error",
            "-show_entries",
            "format=duration",
            "-of",
            "default=nw=1:nk=1",
            str(path),
        ],
        capture_output=True,
        text=True,
        check=True,
    )
    output = completed.stdout.strip()
    if not output:
        raise ExtractionError("ffprobe_no_duration")
    return float(output)


async def extract_audio(
    path: Path,
    *,
    transcriber: Transcriber,
    max_seconds: int | None = None,
) -> str:
    cap = (
        max_seconds
        if max_seconds is not None
        else get_settings().operator_upload_max_audio_seconds
    )
    duration = _probe_duration(path)
    if duration > cap:
        raise ExtractionError("audio_too_long")
    transcript = await asyncio.to_thread(transcriber.transcribe, path, language="ru")
    if not transcript or not transcript.strip():
        raise ExtractionError("empty_text")
    return transcript


async def extract_video(
    path: Path,
    *,
    transcriber: Transcriber,
    max_seconds: int | None = None,
) -> str:
    cap = (
        max_seconds
        if max_seconds is not None
        else get_settings().operator_upload_max_audio_seconds
    )
    duration = _probe_duration(path)
    if duration > cap:
        raise ExtractionError("audio_too_long")
    audio_handle = tempfile.NamedTemporaryFile(suffix=".mp3", delete=False)
    audio_path = Path(audio_handle.name)
    audio_handle.close()
    try:
        subprocess.run(
            [
                "ffmpeg",
                "-y",
                "-i",
                str(path),
                "-vn",
                "-acodec",
                "libmp3lame",
                "-q:a",
                "5",
                str(audio_path),
            ],
            check=True,
            capture_output=True,
        )
        return await extract_audio(audio_path, transcriber=transcriber, max_seconds=cap)
    finally:
        audio_path.unlink(missing_ok=True)


class WhisperTranscriber:
    """Thin wrapper around faster-whisper with lazy model loading."""

    def __init__(self) -> None:
        self._model = None

    def _load(self):  # pragma: no cover - real model load only in production
        from faster_whisper import WhisperModel

        settings = get_settings()
        return WhisperModel(
            settings.faster_whisper_model_size,
            device="cpu",
            compute_type=settings.faster_whisper_compute_type,
            download_root=settings.faster_whisper_cache_dir,
        )

    def transcribe(self, audio_path: Path, *, language: str) -> str:
        if self._model is None:
            self._model = self._load()
        segments, _ = self._model.transcribe(str(audio_path), language=language)
        return " ".join(segment.text.strip() for segment in segments if segment.text)


async def extract_media(
    source_file_type: str,
    path: Path,
    *,
    transcriber: Transcriber,
    max_seconds: int | None = None,
) -> str:
    if source_file_type == "audio":
        return await extract_audio(path, transcriber=transcriber, max_seconds=max_seconds)
    if source_file_type == "video":
        return await extract_video(path, transcriber=transcriber, max_seconds=max_seconds)
    raise ExtractionError(f"unsupported_media_type:{source_file_type}")


def binary_sha256(path: Path) -> str:
    """Streamed SHA-256 of a file's contents (64 KiB blocks)."""
    digest = hashlib.sha256()
    with path.open("rb") as fp:
        while True:
            block = fp.read(65536)
            if not block:
                break
            digest.update(block)
    return digest.hexdigest()
